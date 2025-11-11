import os
import csv
import re
import mimetypes
from django.db import connections
from django.utils import timezone
from PyPDF2 import PdfReader
import docx
from pptx import Presentation

from scraperSite.models import (
    MoodleUrl, MoodleUser,
    MoodleCourseModules, MoodleModules,
    Forum, ForumDiscussion, ForumPost,
    MoodleCourse, MoodleChat, MoodleChatMessage
)

# ----------------------------------------------------
# CONFIG: Moodle dataroot ("moodledata")
# ----------------------------------------------------
FILEDIR_ROOT = r"U:\\"


# ----------------------------------------------------
# URL extraction helpers
# ----------------------------------------------------
# Updated regex: stop before closing brackets, quotes, punctuation
URL_REGEX = re.compile(
    r'https?://[^\s\'"<>)\]}]+',
    re.IGNORECASE
)
TRAILING_CHARS = '.,);:]}>"\''


def clean_url(url):
    """Remove trailing punctuation from extracted URL."""
    return url.rstrip(TRAILING_CHARS)


def extract_urls_from_text(text):
    """Extract all valid URLs from plain text."""
    if not text:
        return []
    found = URL_REGEX.findall(text)
    return [clean_url(u) for u in found]


def extract_urls_from_file(file_path):
    """Extract URLs from different file types (PDF, DOCX, PPTX, TXT)."""
    urls = set()
    if not file_path or not os.path.exists(file_path):
        return []

    mime, _ = mimetypes.guess_type(file_path)

    try:
        if mime == "application/pdf":
            reader = PdfReader(file_path)
            for page in reader.pages:
                text = page.extract_text() or ""
                urls.update(extract_urls_from_text(text))

        elif mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            doc = docx.Document(file_path)
            text = "\n".join(p.text for p in doc.paragraphs)
            urls.update(extract_urls_from_text(text))

        elif mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            pres = Presentation(file_path)
            all_text = []
            for slide in pres.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        all_text.append(shape.text)
            urls.update(extract_urls_from_text("\n".join(all_text)))

        else:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                urls.update(extract_urls_from_text(f.read()))

    except Exception as e:
        print(f"⚠️ Could not extract URLs from file {file_path}: {e}")

    return list(urls)


def write_url_row(writer, url, source, course, author=None):
    """Write one URL record into CSV export."""
    if author:
        author_username = getattr(author, "username", "unknown")
        first = getattr(author, "firstname", "") or ""
        last = getattr(author, "lastname", "") or ""
        author_name = (first + " " + last).strip() or "unknown"
        author_email = getattr(author, "email", "unknown") or "unknown"
    else:
        author_username = author_name = author_email = "unknown"

    writer.writerow([
        url,
        author_username,
        author_name,
        author_email,
        source,
        course.id,
        course.fullname
    ])


# ----------------------------------------------------
# Moodle File API Helpers
# ----------------------------------------------------
def moodle_file_path_from_contenthash(contenthash):
    """Convert Moodle contenthash into a full file path."""
    if not contenthash or len(contenthash) < 4:
        return None
    a, b = contenthash[0:2], contenthash[2:4]
    return os.path.join(FILEDIR_ROOT, a, b, contenthash)


def get_course_module_contextids(course_id, using="moodle"):
    """Return all context IDs for modules in a given Moodle course."""
    sql = """
        SELECT ctx.id
        FROM mdl_course_modules cm
        JOIN mdl_context ctx
          ON ctx.instanceid = cm.id
         AND ctx.contextlevel = 70
        WHERE cm.course = %s
    """
    ids = set()
    with connections[using].cursor() as cur:
        cur.execute(sql, [course_id])
        for (ctx_id,) in cur.fetchall():
            ids.add(ctx_id)
    return ids


def iter_files_for_contextids(context_ids, using="moodle"):
    """Iterate through files linked to Moodle course context IDs."""
    if not context_ids:
        return
    ids = list(context_ids)
    chunk_size = 1000

    sql_base = """
        SELECT contextid, component, filearea, itemid, filename, contenthash, userid
        FROM mdl_files
        WHERE filename <> '.'
          AND contextid IN ({placeholders})
    """

    with connections[using].cursor() as cur:
        for i in range(0, len(ids), chunk_size):
            chunk = ids[i:i+chunk_size]
            placeholders = ",".join(["%s"] * len(chunk))
            sql = sql_base.format(placeholders=placeholders)
            cur.execute(sql, chunk)
            for row in cur.fetchall():
                yield row


def lookup_moodle_user(user_id, using="moodle"):
    """Look up Moodle user safely."""
    if not user_id:
        return None
    return MoodleUser.objects.using(using).filter(id=user_id).first()


# ----------------------------------------------------
# MAIN EXPORT FUNCTION
# ----------------------------------------------------
def export_course_urls(course, export_dir=None, scan_type="auto"):
    """
    Exports all URLs from the course into a CSV-like text file.
    scan_type = "auto" or "manual"  (used for filename prefix)
    """
    today_str = timezone.localtime(timezone.now()).strftime("%Y-%m-%d")
    export_dir = export_dir or os.path.join("url_details", today_str)
    os.makedirs(export_dir, exist_ok=True)

    # 👇 Tag output filename with scan type
    output_csv = os.path.join(export_dir, f"{scan_type}_{course.id}_scanner_input.txt")
    urls_set = set()

    print(f"📁 Exporting URLs for course: {course.fullname}")
    print(f"🗂️ Output file: {output_csv}")

    with open(output_csv, "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f, quoting=csv.QUOTE_ALL)
        writer.writerow(["url", "authorUsername", "authorName", "authorEmail", "source", "courseID", "courseName"])

        # ---------- URL resources ----------
        for u in MoodleUrl.objects.using("moodle").filter(course=course.id):
            author = MoodleUser.objects.using("moodle").filter(id=getattr(u, "userid", None)).first()
            for url in extract_urls_from_text(getattr(u, "externalurl", "") or ""):
                if url not in urls_set:
                    urls_set.add(url)
                    write_url_row(writer, url, "url_resource", course, author)

        # ---------- Forums ----------
        forums = Forum.objects.using("moodle").filter(course=course.id)
        for forum in forums:
            for url in extract_urls_from_text(forum.intro or ""):
                if url not in urls_set:
                    urls_set.add(url)
                    write_url_row(writer, url, "forum_intro", course)

            discussions = ForumDiscussion.objects.using("moodle").filter(forum=forum.id)
            for discussion in discussions:
                author = MoodleUser.objects.using("moodle").filter(id=discussion.userid).first()
                for field in ("name", "content", "message"):
                    text = getattr(discussion, field, "") or ""
                    for url in extract_urls_from_text(text):
                        if url not in urls_set:
                            urls_set.add(url)
                            write_url_row(writer, url, f"forum_discussion_{field}", course, author)

                posts = ForumPost.objects.using("moodle").filter(discussion=discussion.id)
                for post in posts:
                    post_author = MoodleUser.objects.using("moodle").filter(id=post.userid).first()
                    for url in extract_urls_from_text(post.message or ""):
                        if url not in urls_set:
                            urls_set.add(url)
                            write_url_row(writer, url, "forum_post", course, post_author)

        # ---------- Chats ----------
        chats = MoodleChat.objects.using("moodle").filter(course=course.id)
        for chat in chats:
            for url in extract_urls_from_text(chat.intro or ""):
                if url not in urls_set:
                    urls_set.add(url)
                    write_url_row(writer, url, "chat_intro", course)

            msgs = MoodleChatMessage.objects.using("moodle").filter(chatid=chat.id)
            for msg in msgs:
                msg_author = MoodleUser.objects.using("moodle").filter(id=msg.userid).first()
                for url in extract_urls_from_text(msg.message or ""):
                    if url not in urls_set:
                        urls_set.add(url)
                        write_url_row(writer, url, "chat_message", course, msg_author)

        # ---------- Moodle Files ----------
        context_ids = get_course_module_contextids(course.id, "moodle")
        for contextid, component, filearea, itemid, filename, contenthash, file_userid in iter_files_for_contextids(context_ids):
            file_path = moodle_file_path_from_contenthash(contenthash)
            if not file_path or not os.path.exists(file_path):
                continue

            file_urls = extract_urls_from_file(file_path)
            if not file_urls:
                continue

            author = lookup_moodle_user(file_userid)
            source = f"{component}:{filearea}"

            for url in file_urls:
                if url not in urls_set:
                    urls_set.add(url)
                    write_url_row(writer, url, source, course, author)

    print(f"✅ EXPORT COMPLETE — Total URLs: {len(urls_set)}")
    print(f"📄 Saved to: {output_csv}")
    return output_csv
